"""
سكريبت لإنشاء ملف PowerPoint من محتوى مشروع عَقِلْها
يتطلب تثبيت: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # إنشاء عرض تقديمي جديد
    prs = Presentation()
    
    # الألوان من القالب (Dark Olive Green)
    bg_color = RGBColor(85, 107, 47)  # Dark Olive Green
    title_color = RGBColor(255, 165, 0)  # Orange
    text_color = RGBColor(255, 255, 255)  # White
    
    # الشريحة 1: أعضاء الفريق
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title1.text_frame
    title_frame.text = "أبـشـر طويق - قالب تحكيم المشاريع"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = title_color
    title_frame.paragraphs[0].font.bold = True
    
    # جدول أعضاء الفريق
    members_text = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    members_frame = members_text.text_frame
    members_frame.text = "أعضاء الفريق\n\n| 4 | 5 | 6 | 4 |\n| اسم العضو | اسم العضو | اسم العضو | اسم العضو |"
    members_frame.paragraphs[0].font.size = Pt(18)
    members_frame.paragraphs[0].font.color.rgb = text_color
    
    # الشريحة 2: المشكلة وحلها
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title2_frame = title2.text_frame
    title2_frame.text = "المشكلة وحلّها"
    title2_frame.paragraphs[0].font.size = Pt(36)
    title2_frame.paragraphs[0].font.color.rgb = title_color
    title2_frame.paragraphs[0].font.bold = True
    
    # المشكلة
    problem_text = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    problem_frame = problem_text.text_frame
    problem_frame.text = """المشكلة:

• ازدحام مروري مزمن في المدن الرئيسية
• نقص المعلومات الفورية عن حالة الطرق
• عدم وجود تنبؤات ذكية بالازدحام
• صعوبة اتخاذ قرارات سريعة
• التأثير السلبي على الوقت والوقود"""
    problem_frame.paragraphs[0].font.size = Pt(16)
    problem_frame.paragraphs[0].font.color.rgb = text_color
    
    # الحل
    solution_text = slide2.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5))
    solution_frame = solution_text.text_frame
    solution_frame.text = """الحل:

عَقِلْها - نظام ذكي شامل

• استخدام التكنولوجيا الحديثة
• الاستجابة لاحتياجات العملاء
• الاستفادة من الاتجاهات
• التوقيت الذكي"""
    solution_frame.paragraphs[0].font.size = Pt(16)
    solution_frame.paragraphs[0].font.color.rgb = text_color
    
    # الشريحة 3: البيانات المستخدمة
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title3_frame = title3.text_frame
    title3_frame.text = "البيانات المستخدمة"
    title3_frame.paragraphs[0].font.size = Pt(36)
    title3_frame.paragraphs[0].font.color.rgb = title_color
    title3_frame.paragraphs[0].font.bold = True
    
    data_text = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    data_frame = data_text.text_frame
    data_frame.text = """مصادر البيانات:

1. Google Maps Platform (40%)
   - Google Traffic Layer
   - Google Directions API
   - Google Places API

2. OpenWeatherMap API (30%)
   - بيانات الطقس الحالية
   - تنبؤات لمدة 16 يوم

3. قاعدة البيانات المحلية (20%)
   - بيانات تاريخية
   - سجلات التنبيهات

4. بيانات المستخدمين (10%)
   - مواقع مجهولة الهوية
   - تفضيلات المسارات"""
    data_frame.paragraphs[0].font.size = Pt(16)
    data_frame.paragraphs[0].font.color.rgb = text_color
    
    # الشريحة 4: التقنيات المستخدمة
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title4_frame = title4.text_frame
    title4_frame.text = "التقنيات المستخدمة"
    title4_frame.paragraphs[0].font.size = Pt(36)
    title4_frame.paragraphs[0].font.color.rgb = title_color
    title4_frame.paragraphs[0].font.bold = True
    
    tech_text = slide4.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tech_frame = tech_text.text_frame
    tech_frame.text = """Frontend:
• Next.js 14 - React Framework
• Tailwind CSS - التصميم
• Google Maps JavaScript API
• TypeScript

Backend:
• Next.js API Routes
• Prisma ORM
• MySQL Database
• React Query

APIs:
• Google Maps Platform
• OpenWeatherMap API
• Google Routes API"""
    tech_frame.paragraphs[0].font.size = Pt(16)
    tech_frame.paragraphs[0].font.color.rgb = text_color
    
    # الشريحة 5: وصف الفكرة
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title5_frame = title5.text_frame
    title5_frame.text = "وصف الفكرة"
    title5_frame.paragraphs[0].font.size = Pt(36)
    title5_frame.paragraphs[0].font.color.rgb = title_color
    title5_frame.paragraphs[0].font.bold = True
    
    idea_text = slide5.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    idea_frame = idea_text.text_frame
    idea_frame.text = """عنوان: عَقِلْها - نظام تحليل الازدحام المروري الذكي

الفوائد:
• توفير الوقت والوقود
• تحسين السلامة
• التنبؤ الذكي
• واجهة سهلة الاستخدام

الابتكارات:
• تكامل كامل مع Google Maps
• نظام تحذيرات شامل (12 نوع)
• تحليل الطقس لمدة 16 يوم
• مسارات بديلة محسّنة"""
    idea_frame.paragraphs[0].font.size = Pt(16)
    idea_frame.paragraphs[0].font.color.rgb = text_color
    
    # الشريحة 6: كيفية توفير البيانات واستخدامها
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title6_frame = title6.text_frame
    title6_frame.text = "كيفية توفير هذه البيانات وكيفية استخدامها"
    title6_frame.paragraphs[0].font.size = Pt(28)
    title6_frame.paragraphs[0].font.color.rgb = title_color
    title6_frame.paragraphs[0].font.bold = True
    
    data_usage_text = slide6.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    data_usage_frame = data_usage_text.text_frame
    data_usage_frame.text = """كيفية الحصول على البيانات:
• Google Maps API - مفتاح API من Google Cloud
• OpenWeatherMap API - التسجيل والحصول على Key
• قاعدة البيانات المحلية - استخدام Prisma ORM
• بيانات المستخدمين - جمع مجهول الهوية

كيفية استخدام البيانات:
• التحليل الفوري - عرض حالة الازدحام
• التنبؤ الذكي - تحليل الأنماط
• التنبيهات - تحليل المخاطر
• التحسين المستمر - تحسين الخوارزميات"""
    data_usage_frame.paragraphs[0].font.size = Pt(14)
    data_usage_frame.paragraphs[0].font.color.rgb = text_color
    
    # الشريحة 7: مواءمة الفكرة
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title7_frame = title7.text_frame
    title7_frame.text = "مواءمة الفكرة"
    title7_frame.paragraphs[0].font.size = Pt(36)
    title7_frame.paragraphs[0].font.color.rgb = title_color
    title7_frame.paragraphs[0].font.bold = True
    
    alignment_text = slide7.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    alignment_frame = alignment_text.text_frame
    alignment_frame.text = """1. تحسين الأداء:
   • تقليل وقت السفر بنسبة 40%
   • توفير الوقود والطاقة

2. تعزيز السلامة:
   • تنبيهات عن الحوادث
   • تحذيرات الطقس
   • تحليل الرؤية السيئة

3. تحسين تفاعل المستخدمين:
   • واجهة سهلة الاستخدام
   • تنبيهات مخصصة
   • معلومات واضحة

4. الابتكار التكنولوجي:
   • استخدام الذكاء الاصطناعي
   • تكامل مع APIs حديثة"""
    alignment_frame.paragraphs[0].font.size = Pt(16)
    alignment_frame.paragraphs[0].font.color.rgb = text_color
    
    # الشريحة 8: ملخص
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    title8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title8_frame = title8.text_frame
    title8_frame.text = "ملخص"
    title8_frame.paragraphs[0].font.size = Pt(36)
    title8_frame.paragraphs[0].font.color.rgb = title_color
    title8_frame.paragraphs[0].font.bold = True
    
    summary_text = slide8.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    summary_frame = summary_text.text_frame
    summary_frame.text = """الأهداف:
• حل مشكلة الازدحام المروري
• توفير الوقت والوقود
• تحسين السلامة على الطرق

المخرجات:
• نظام ويب متكامل
• واجهة سهلة الاستخدام
• نظام تنبيهات ذكي
• تحليل فوري للازدحام

النتائج:
• تكامل كامل مع Google Maps
• نظام تحذيرات شامل
• تحليل الطقس المتقدم
• مسارات بديلة محسّنة"""
    summary_frame.paragraphs[0].font.size = Pt(16)
    summary_frame.paragraphs[0].font.color.rgb = text_color
    
    # الشريحة 9: الاختبار/التحقق
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title9_frame = title9.text_frame
    title9_frame.text = "الاختبار/التحقق"
    title9_frame.paragraphs[0].font.size = Pt(36)
    title9_frame.paragraphs[0].font.color.rgb = title_color
    title9_frame.paragraphs[0].font.bold = True
    
    test_text = slide9.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    test_frame = test_text.text_frame
    test_frame.text = """ما تم إنجازه:

✅ نظام الخرائط:
   • تكامل مع Google Maps Traffic Layer
   • عرض حالة الازدحام في الوقت الفعلي

✅ نظام التنبيهات:
   • 12 نوع تحذير مختلف
   • تنبيهات فورية

✅ تحليل الطقس:
   • تحليل شامل لمدة 16 يوم
   • تحديد الطرق الآمنة

✅ التنبؤات:
   • تنبؤات لمدة 60 دقيقة
   • تحليل الأنماط التاريخية"""
    test_frame.paragraphs[0].font.size = Pt(16)
    test_frame.paragraphs[0].font.color.rgb = text_color
    
    # الشريحة 10: العرض التوضيحي
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    title10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title10_frame = title10.text_frame
    title10_frame.text = "العرض التوضيحي"
    title10_frame.paragraphs[0].font.size = Pt(36)
    title10_frame.paragraphs[0].font.color.rgb = title_color
    title10_frame.paragraphs[0].font.bold = True
    
    demo_text = slide10.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    demo_frame = demo_text.text_frame
    demo_frame.text = """اللقطات والشرائح التوضيحية:

1. شاشة الخريطة الرئيسية
   • عرض الخريطة مع حالة الازدحام
   • الألوان المعبرة

2. شاشة التنبيهات
   • قائمة التنبيهات النشطة
   • تفاصيل كل تنبيه

3. شاشة تحليل الطقس
   • عرض الطرق الآمنة وغير الآمنة
   • تحذيرات الطقس

4. شاشة المسارات
   • حساب المسار الأسرع
   • عرض المسافة والوقت

ملاحظة: تم إنجاز أكثر من 30% من المشروع"""
    demo_frame.paragraphs[0].font.size = Pt(16)
    demo_frame.paragraphs[0].font.color.rgb = text_color
    
    # الشريحة 11: التحديات والخطط المستقبلية
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    title11 = slide11.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title11_frame = title11.text_frame
    title11_frame.text = "التحديات والخطط المستقبلية"
    title11_frame.paragraphs[0].font.size = Pt(28)
    title11_frame.paragraphs[0].font.color.rgb = title_color
    title11_frame.paragraphs[0].font.bold = True
    
    challenges_text = slide11.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    challenges_frame = challenges_text.text_frame
    challenges_frame.text = """التحديات:
• تكامل مع APIs متعددة
• معالجة البيانات الكبيرة
• ضمان سرعة الاستجابة

ما تحتاج المساعدة فيه:
• مفاتيح API إضافية
• خوادم أقوى
• إرشاد من الخبراء

العمل المستقبلي (70%):
• تحسين دقة التنبؤات
• إضافة دعم لمدن إضافية
• تحسين سرعة الاستجابة
• إطلاق النسخة النهائية"""
    challenges_frame.paragraphs[0].font.size = Pt(14)
    challenges_frame.paragraphs[0].font.color.rgb = text_color
    
    # الشريحة 12: Timeline
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    title12 = slide12.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title12_frame = title12.text_frame
    title12_frame.text = "Timeline"
    title12_frame.paragraphs[0].font.size = Pt(36)
    title12_frame.paragraphs[0].font.color.rgb = title_color
    title12_frame.paragraphs[0].font.bold = True
    
    timeline_text = slide12.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    timeline_frame = timeline_text.text_frame
    timeline_frame.text = """الشهر الأول:
• الأسبوع 1-2: التخطيط والتصميم
• الأسبوع 3-4: التطوير الأساسي

الشهر الثاني:
• الأسبوع 1-2: التكامل مع APIs
• الأسبوع 3-4: نظام التنبيهات

الشهر الثالث:
• الأسبوع 1-2: تحليل الطقس
• الأسبوع 3-4: التحسينات والاختبار

الشهر الرابع:
• الأسبوع 1-2: التحسينات النهائية
• الأسبوع 3-4: الإطلاق والتحسينات

الحالة الحالية: ✅ أكثر من 30% مكتمل"""
    timeline_frame.paragraphs[0].font.size = Pt(16)
    timeline_frame.paragraphs[0].font.color.rgb = text_color
    
    # الشريحة 13: Thank you
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    title13 = slide13.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(2))
    title13_frame = title13.text_frame
    title13_frame.text = "شكراً لكم"
    title13_frame.paragraphs[0].font.size = Pt(48)
    title13_frame.paragraphs[0].font.color.rgb = title_color
    title13_frame.paragraphs[0].font.bold = True
    title13_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle13 = slide13.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1))
    subtitle13_frame = subtitle13.text_frame
    subtitle13_frame.text = "عَقِلْها\nنظام تحليل الازدحام المروري الذكي"
    subtitle13_frame.paragraphs[0].font.size = Pt(24)
    subtitle13_frame.paragraphs[0].font.color.rgb = text_color
    subtitle13_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # حفظ الملف
    prs.save('عرض_عقيلها.pptx')
    print("✅ تم إنشاء ملف PowerPoint بنجاح: عرض_عقيلها.pptx")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ خطأ: يجب تثبيت مكتبة python-pptx أولاً")
        print("قم بتشغيل: pip install python-pptx")
    except Exception as e:
        print(f"❌ حدث خطأ: {e}")

